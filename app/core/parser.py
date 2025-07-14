import base64
import io
import tempfile
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Callable, Dict, List, Tuple

import aspose.words as aw
import cchardet
import fitz
import mammoth
import openpyxl
import pandas as pd
import pdfplumber
import requests
from docx import Document
from langchain_community.document_loaders import UnstructuredEPubLoader
from loguru import logger
from markdownify import MarkdownConverter
from pandas import DataFrame
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PyPDF2 import PdfReader

from app.middlewares.exception.common import CommonException
from app.schemas.file import (AllowedAudioSuffix, AllowedDataBaseSuffix,
                              AllowedDocumentSuffix, AllowedFileSuffix,
                              AllowedImageSuffix, AllowedTableSuffix,
                              AllowedVideoSuffix, MarkdownAssetContent)
from app.schemas.public import AssetInfo, AssetType
from app.services.external.baidu import BaiduOcr
from app.services.external.iflytek import IflytekSpeech
from app.services.external.llama import Llama

ImageSaver = Callable[[io.BytesIO, str, str], Tuple[bool, str, str]]

PPT_PAGE_SPLITER = '# <第{page}页>\n\n'


def detect_encoding(file_bytes: io.BytesIO) -> str:
    result = cchardet.detect(file_bytes.getvalue())
    logger.info(f'detected encoding: {result}')
    file_bytes.seek(0)
    return result.get('encoding', 'utf-8') or 'utf-8'


class File2MarkdownBaseParser:

    def __init__(self, file_bytes: io.BytesIO, image_save_fn: ImageSaver = None):
        self.file_bytes = file_bytes
        self.tables: List[DataFrame] = []
        self.assets: List[AssetInfo] = []
        self.path_to_assets: Dict[str, AssetInfo] = None
        self.image_save_fn = image_save_fn

    def save_images(self, file_bytes: io.BytesIO, name: str, suffix: str) -> Tuple[bool, str, str]:
        if self.image_save_fn is not None:
            is_success, path, aid = self.image_save_fn(file_bytes, name, suffix)
            return is_success, path, aid
        else:
            file_name = f'{name}{suffix}'
            with open(path := Path().absolute() / f'{name}{suffix}', 'wb') as f:
                f.write(file_bytes.getvalue())
            return True, str(path), file_name

    def get_asset_by_path(self, path: str) -> AssetInfo | None:
        if not self.path_to_assets:
            self.path_to_assets = {}
            for asset in self.assets:
                self.path_to_assets[asset.path] = asset
        return self.path_to_assets.get(path, None)

    def load(self) -> str | DataFrame:
        return ''


class Text2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO):
        super().__init__(file_bytes)

    def load(self) -> str:
        return self.file_bytes.getvalue().decode(detect_encoding(self.file_bytes))


class DataBase2MarkdownParser(Text2MarkdownParser):
    ...


class Epub2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO):
        super().__init__(file_bytes)

    def load(self) -> str:
        with tempfile.NamedTemporaryFile() as temp_file:
            file_name = f'{temp_file.name}.epub'
            with open(file_name, 'wb') as f:
                f.write(self.file_bytes.getvalue())
            return UnstructuredEPubLoader(file_name).load()[0].page_content


class Word2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO, suffix: str, image_save_fn: ImageSaver):
        super().__init__(file_bytes, image_save_fn)
        self.suffix = suffix
        self.md_converter = MarkdownConverter(heading_style='')

    def get_img_convert_fn(self):

        def convert_img(*args):
            src = args[0].get('src', '')
            if ',' not in src:
                return ''

            meta, base64_str = src.split(',')
            if 'image/png' in meta:
                suffix = '.png'
            elif 'image/jpeg' in meta:
                suffix = '.jpeg'
            elif 'image/jpg' in meta:
                suffix = '.jpg'
            else:
                return ''

            data = io.BytesIO(base64.b64decode(base64_str))
            name = f'图片{len(self.assets) + 1}'
            is_success, path, key = self.save_images(data, name, suffix)
            if is_success:
                asset_info = AssetInfo(aid=int(key),
                                       path=path,
                                       atype=AssetType.IMAGE,
                                       size=data.getbuffer().nbytes,
                                       suffix=suffix,
                                       title=name)
                self.assets.append(asset_info)
                return '\n\n' + asset_info.to_md_content() + '\n\n'
            else:
                return ''

        return convert_img

    def get_table_convert_fn(self):

        def convert_table(*args):
            try:
                dfs = pd.read_html(io.StringIO(str(args[0])))
            except Exception as e:
                logger.error(f'Convert table error: {e}')
                dfs = []
            if len(dfs) == 0:
                return ''
            else:
                self.tables.extend(dfs)
                return '\n\n' + args[1] + '\n\n'

        return convert_table

    @staticmethod
    def convert_doc_bytes(file_bytes: io.BytesIO):
        logger.info('Converting doc to docx')
        docx_file_bytes = io.BytesIO()
        aw.Document(file_bytes).save(docx_file_bytes, aw.SaveFormat.DOCX)
        docx_file_bytes.seek(0)
        doc = Document(docx_file_bytes)

        # 去除开头水印
        fp = doc.paragraphs[0]._element
        fp.getparent().remove(fp)
        fp.element = fp._p = None
        byte_io = io.BytesIO()

        # 结尾水印
        lp = doc.paragraphs[-1]._element
        if 'This document was truncated here because it was created in the Evaluation Mode.' in lp.text:
            lp.getparent().remove(lp)
            lp.element = lp._p = None

        doc.save(byte_io)
        byte_io.seek(0)
        return byte_io

    def extract_body(self) -> str:
        b = self.convert_doc_bytes(self.file_bytes) if self.suffix == '.doc' else self.file_bytes
        html = mammoth.convert_to_html(b).value
        self.md_converter.convert_img = self.get_img_convert_fn()
        self.md_converter.convert_table = self.get_table_convert_fn()
        markdown = self.md_converter.convert(html).replace('\\', '')
        return markdown

    def load(self) -> str:
        return self.extract_body()


class Table2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO, suffix: str, image_save_fn: ImageSaver):
        super().__init__(file_bytes, image_save_fn)
        self.suffix = suffix
        self.images = []
        self.table = None

    def extract_table(self):
        if self.suffix == '.xlsx':
            df = pd.read_excel(self.file_bytes, engine='openpyxl')
        elif self.suffix == '.csv':
            df = pd.read_csv(self.file_bytes, encoding=detect_encoding(self.file_bytes))
        elif self.suffix == '.xls':
            df = pd.read_excel(self.file_bytes)
        else:
            raise Exception('不支持该格式')
        self.table = df.fillna('')
        return self.table

    def extract_images(self):
        if self.suffix == '.xlsx':
            wb = openpyxl.load_workbook(self.file_bytes)
            sheet = wb[wb.sheetnames[0]]
            self.images = sheet._images
        return self.images

    def merge_text_images(self):

        for img_idx, img in enumerate(self.images):
            pos = img.anchor._from
            row, col = pos.row - 1, pos.col
            file_bytes = io.BytesIO(img._data())
            suffix = '.png'
            name = f'图片{img_idx + 1}'
            is_success, path, key = self.save_images(file_bytes, name, suffix)
            if is_success:
                asset = AssetInfo(aid=int(key),
                                  path=path,
                                  atype=AssetType.IMAGE,
                                  size=file_bytes.getbuffer().nbytes,
                                  suffix=suffix,
                                  title=name)
                self.table.iloc[row, col] = self.table.iloc[row, col] + '\n\n' + asset.to_md_content()

        return self.table

    def load(self) -> DataFrame:
        self.extract_table()
        self.extract_images()
        self.merge_text_images()
        return self.table.to_markdown(index=False)


class Pdf2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO, image_save_fn: ImageSaver):
        super().__init__(file_bytes, image_save_fn)
        self.is_ppt = self.is_from_ppt()
        self.page2assets = defaultdict(list)
        self.page2tables = defaultdict(list)
        self.page2text = defaultdict(str)

    def is_from_ppt(self):
        creator = PdfReader(self.file_bytes).metadata.creator
        logger.info(f'PDF creator: {creator}')
        creator = str(creator)

        if 'PowerPoint' in creator:
            return True
        if 'WPS 演示' in creator:
            return True

        return False

    def get_doc(self):
        with tempfile.NamedTemporaryFile() as temp_file:
            file_name = f'{temp_file.name}.pdf'
            with open(file_name, 'wb') as f:
                f.write(self.file_bytes.getvalue())
            doc = Llama().parse_pdf(file_name)
            return doc

    def get_page_doc(self):
        doc = fitz.open(stream=self.file_bytes, filetype='pdf')
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            self.page2text[page_num] = page.get_text()
        return doc.page_count

    def extract_image(self):

        with pdfplumber.open(self.file_bytes) as pdf:
            img_idx = 1
            for page_num, page in enumerate(pdf.pages):
                for idx, image in enumerate(page.images, start=1):
                    img = page.within_bbox((image['x0'], image['top'], image['x1'], image['bottom'])).to_image()
                    file_bytes = io.BytesIO()
                    img.save(file_bytes, format='PNG')
                    file_bytes.seek(0)
                    suffix = '.png'
                    name = f'图片{img_idx}'
                    is_success, path, key = self.save_images(file_bytes, name, suffix)
                    if is_success:
                        asset = AssetInfo(aid=int(key),
                                          path=path,
                                          atype=AssetType.IMAGE,
                                          size=file_bytes.getbuffer().nbytes,
                                          suffix=suffix,
                                          title=name)
                        self.assets.append(asset)
                        self.page2assets[page_num].append(asset)

                    img_idx += 1

        return self.assets

    def extract_tables(self):

        with pdfplumber.open(self.file_bytes) as pdf:
            for idx, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table_number, table in enumerate(tables):
                    df = pd.DataFrame(table)
                    self.tables.append(df)
                    self.page2tables[idx].append(df)
        return self.tables

    def load_for_full(self):
        result = self.get_doc().text
        assets = self.extract_image()
        if len(assets) > 0:
            result += '\n\n图片：\n\n' + '\n\n'.join([i.to_md_content() for i in assets])
        if len(self.tables) > 0:
            result += '\n\n表格：\n\n' + '\n\n'.join([x.to_markdown(index=False) for x in self.tables])
        return result

    def load_for_pages(self):
        total_page = self.get_page_doc()
        self.extract_image()
        self.extract_tables()

        content_ls = []
        for idx in range(total_page):
            content = PPT_PAGE_SPLITER.format(page=idx + 1)

            if len(text := self.page2text.get(idx, '')) > 0:
                content += text + '\n\n'
            if len(tables := self.page2tables.get(idx, [])) > 0:
                content += '\n\n'.join([x.to_markdown(index=False) for x in tables]) + '\n\n'
            if len(assets := self.page2assets.get(idx, [])) > 0:
                content += '\n\n'.join([x.to_md_content() for x in assets]) + '\n\n'

            content_ls.append(content)

        return '\n\n'.join(content_ls)

    def load(self) -> str:
        if self.is_ppt:
            return self.load_for_pages()
        else:
            return self.load_for_full()


class Ppt2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO, image_save_fn: ImageSaver):
        super().__init__(file_bytes, image_save_fn)
        self.img_idx = 1
        self.slides = []

    def extract_slide(self, slide):

        slide_data = {'text': [], 'image': [], 'table': []}

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content = [run.text for run in paragraph.runs]
                    slide_data['text'].append(''.join(content))

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                data = [[cell.text for cell in row.cells] for row in table.rows]
                slide_data['table'].append(pd.DataFrame(data[1:], columns=data[0]))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                _suffix = image.content_type.split('/')[-1]
                suffix = f'.{_suffix}'
                name = f'图片{self.img_idx}'
                file_bytes = io.BytesIO(image.blob)
                is_success, path, key = self.save_images(file_bytes, name, suffix)
                if is_success:
                    asset = AssetInfo(aid=int(key), path=path, atype=AssetType.IMAGE,
                                      size=file_bytes.getbuffer().nbytes,
                                      suffix=f'.{suffix}', title=name)
                    slide_data['image'].append(asset)
                    self.assets.append(asset)
                self.img_idx += 1

        self.slides.append(slide_data)

    def load(self) -> str:
        pptx = Presentation(self.file_bytes)
        for slide in pptx.slides:
            self.extract_slide(slide)

        content = []
        for idx, data in enumerate(self.slides):
            text = PPT_PAGE_SPLITER.format(page=idx + 1)
            if len(data['text']) > 0:
                text += '* 文本内容：\n\n' + '\n\n'.join(data['text']) + '\n\n'
            if len(data['table']) > 0:
                text += '* 表格内容：\n\n' + '\n\n'.join([x.to_markdown(index=False) for x in data['table']]) + '\n\n'
            if len(data['image']) > 0:
                text += '* 图片：\n\n' + '\n\n'.join([x.to_md_content() for x in data['image']])
            content.append(text)
        return '\n\n'.join(content)


class Image2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO):
        super().__init__(file_bytes)

    def load(self) -> str:
        client = BaiduOcr()
        result = client.accurate(self.file_bytes.getvalue())
        data = BaiduOcr.parse_simple_paragraph(result)
        return data


class Media2MarkdownParser(File2MarkdownBaseParser):

    def __init__(self, file_bytes: io.BytesIO, suffix: str):
        super().__init__(file_bytes)
        self.suffix = suffix

    def load(self) -> str:
        data = IflytekSpeech(file_bytes=self.file_bytes, suffix=self.suffix).get_result()
        result = IflytekSpeech.parse_text(data)
        return result


class Zip2MarkdownParser(File2MarkdownBaseParser):
    def __init__(self, file_bytes: io.BytesIO,  image_save_fn: ImageSaver):
        super().__init__(file_bytes, image_save_fn)
        self.zip_path: Path = None

    def get_md_string(self) -> str:
        md_file = [x for x in self.zip_path.iterdir() if x.is_file() and x.suffix.endswith('.md')]
        if len(md_file) != 1:
            raise Exception('找不到 .md 文件')

        with open(md_file[0], 'rb') as f:
            md_file_bytes = io.BytesIO(f.read())

        md_string = Text2MarkdownParser(file_bytes=md_file_bytes).load()
        return md_string

    def convert_images(self, md_string) -> str:

        md_replace = md_string
        for i in MarkdownAssetContent.parse_from_text(md_string):
            replace_image_str = ''
            if not str(i.path).startswith(('https://', 'http://')):
                image_file_path = self.zip_path / i.path
                if image_file_path.exists():
                    logger.info(f'Found image: {image_file_path}')
                    with open(image_file_path, 'rb') as f:
                        data = io.BytesIO(f.read())

                    name = f'图片{len(self.assets) + 1}'
                    suffix = image_file_path.suffix
                    is_success, path, key = self.save_images(data, name, suffix)
                    if is_success:
                        asset_info = AssetInfo(aid=int(key),
                                               path=path,
                                               atype=AssetType.IMAGE,
                                               size=data.getbuffer().nbytes,
                                               suffix=suffix,
                                               title=name)
                        self.assets.append(asset_info)
                        replace_image_str = asset_info.to_md_content()
            else:
                try:
                    response = requests.get(str(i.path))
                    if response.status_code != 200:
                        raise Exception(f'Failed to download image: {i.path}, status_code: {response.status_code}')

                    data = io.BytesIO(response.content)
                    name = f'图片{len(self.assets) + 1}'
                    suffix = f".{str(i.path).split('.')[-1]}"
                    is_success, path, key = self.save_images(data, name, suffix)
                    if is_success:
                        asset_info = AssetInfo(aid=int(key),
                                               path=path,
                                               atype=AssetType.IMAGE,
                                               size=data.getbuffer().nbytes,
                                               suffix=suffix,
                                               title=name)
                        self.assets.append(asset_info)
                        replace_image_str = asset_info.to_md_content()

                except Exception as e:
                    logger.error(f'Failed to process url image: {i.path}, error: {e}')

            md_replace = md_replace.replace(i.raw, replace_image_str)

        return md_replace

    def load(self) -> str:
        with tempfile.TemporaryDirectory() as temp_dir:
            self.zip_path = Path(temp_dir)
            with zipfile.ZipFile(self.file_bytes, 'r') as zip_ref:
                zip_ref.extractall(str(self.zip_path))
                md_string = self.get_md_string()
                md_string = self.convert_images(md_string=md_string)
        return md_string


class MarkdownParserFactory:

    @staticmethod
    def get_parser(file_bytes: io.BytesIO,
                   suffix: AllowedFileSuffix,
                   image_save_fn: ImageSaver
                   ) -> File2MarkdownBaseParser:
        match suffix:
            case AllowedDocumentSuffix.MD | AllowedDocumentSuffix.TXT:
                parser = Text2MarkdownParser(file_bytes=file_bytes)
            case AllowedDocumentSuffix.ZIP:
                parser = Zip2MarkdownParser(file_bytes=file_bytes, image_save_fn=image_save_fn)
            case AllowedDataBaseSuffix():
                parser = DataBase2MarkdownParser(file_bytes=file_bytes)
            case AllowedDocumentSuffix.DOC | AllowedDocumentSuffix.DOCX:
                parser = Word2MarkdownParser(file_bytes=file_bytes, suffix=suffix.value, image_save_fn=image_save_fn)
            case AllowedDocumentSuffix.PDF:
                parser = Pdf2MarkdownParser(file_bytes=file_bytes, image_save_fn=image_save_fn)
            case AllowedDocumentSuffix.PPTX:
                parser = Ppt2MarkdownParser(file_bytes=file_bytes, image_save_fn=image_save_fn)
            case AllowedTableSuffix.CSV | AllowedTableSuffix.XLSX | AllowedTableSuffix.XLS:
                parser = Table2MarkdownParser(file_bytes=file_bytes, suffix=suffix.value, image_save_fn=image_save_fn)
            case AllowedDocumentSuffix.EPUB:
                parser = Epub2MarkdownParser(file_bytes=file_bytes)
            case AllowedImageSuffix():
                parser = Image2MarkdownParser(file_bytes=file_bytes)
            case AllowedAudioSuffix() | AllowedVideoSuffix():
                parser = Media2MarkdownParser(file_bytes=file_bytes, suffix=suffix.value)
            case _:
                raise CommonException.system_error(f'Unsupported file suffix: {suffix}')

        logger.info(f'Using `{parser.__class__}` to parse markdown')
        return parser
